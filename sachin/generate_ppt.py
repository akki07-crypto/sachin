import sys
import subprocess

# Auto-install python-pptx if not installed
try:
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.dml.color import RGBColor
    from pptx.enum.shapes import MSO_SHAPE
except ImportError:
    print("python-pptx not found. Installing now...")
    subprocess.check_call([sys.executable, "-m", "pip", "install", "python-pptx"])
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.dml.color import RGBColor
    from pptx.enum.shapes import MSO_SHAPE

def create_presentation():
    prs = Presentation()
    # Set to 16:9 widescreen
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    # Color Palette (Tech. Point Theme)
    DARK_BG = RGBColor(7, 9, 19)
    CYAN_PRIMARY = RGBColor(0, 242, 254)
    MAGENTA_SECONDARY = RGBColor(243, 85, 136)
    WHITE_TEXT = RGBColor(255, 255, 255)
    MUTED_TEXT = RGBColor(176, 181, 198)
    CARD_BG = RGBColor(20, 24, 38)

    # Helper: Set slide background to dark
    def set_dark_background(slide):
        bg = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(7.5)
        )
        bg.fill.solid()
        bg.fill.fore_color.rgb = DARK_BG
        bg.line.fill.background()
        return bg

    # Helper: Add Slide Header Title
    def add_slide_header(slide, title_text):
        tx_box = slide.shapes.add_textbox(Inches(0.8), Inches(0.6), Inches(11.7), Inches(1.0))
        tf = tx_box.text_frame
        tf.word_wrap = True
        tf.margin_left = tf.margin_top = tf.margin_right = tf.margin_bottom = 0
        p = tf.paragraphs[0]
        p.text = title_text
        p.font.name = 'Arial'
        p.font.size = Pt(36)
        p.font.bold = True
        p.font.color.rgb = CYAN_PRIMARY
        
        # Add accent line under header
        line = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, Inches(0.8), Inches(1.5), Inches(1.5), Inches(0.06)
        )
        line.fill.solid()
        line.fill.fore_color.rgb = MAGENTA_SECONDARY
        line.line.fill.background()

    # ----------------------------------------------------
    # SLIDE 1: Title Slide (Dark Cover)
    # ----------------------------------------------------
    slide_layout = prs.slide_layouts[6] # Blank
    slide1 = prs.slides.add_slide(slide_layout)
    set_dark_background(slide1)

    # Add design shapes
    decor = slide1.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0.8), Inches(3.2), Inches(0.12), Inches(2.2)
    )
    decor.fill.solid()
    decor.fill.fore_color.rgb = CYAN_PRIMARY
    decor.line.fill.background()

    title_box = slide1.shapes.add_textbox(Inches(1.2), Inches(3.0), Inches(11.0), Inches(2.5))
    tf1 = title_box.text_frame
    tf1.word_wrap = True
    tf1.margin_left = tf1.margin_top = tf1.margin_right = tf1.margin_bottom = 0

    p1 = tf1.paragraphs[0]
    p1.text = "TECH. POINT ACADEMY"
    p1.font.name = 'Arial'
    p1.font.size = Pt(56)
    p1.font.bold = True
    p1.font.color.rgb = WHITE_TEXT

    p2 = tf1.add_paragraph()
    p2.text = "Modern Animated Web Platform & Secure Backend"
    p2.font.name = 'Arial'
    p2.font.size = Pt(22)
    p2.font.color.rgb = CYAN_PRIMARY
    p2.space_before = Pt(10)

    p3 = tf1.add_paragraph()
    p3.text = "Coding Your Thoughts Into Reality"
    p3.font.name = 'Arial'
    p3.font.size = Pt(16)
    p3.font.italic = True
    p3.font.color.rgb = MUTED_TEXT
    p3.space_before = Pt(20)

    # ----------------------------------------------------
    # SLIDE 2: Project Overview & Objectives
    # ----------------------------------------------------
    slide2 = prs.slides.add_slide(slide_layout)
    set_dark_background(slide2)
    add_slide_header(slide2, "Project Overview & Objectives")

    # Content Box
    overview_box = slide2.shapes.add_textbox(Inches(0.8), Inches(2.2), Inches(11.7), Inches(4.5))
    tf2 = overview_box.text_frame
    tf2.word_wrap = True
    
    p_intro = tf2.paragraphs[0]
    p_intro.text = "Tech. Point is a premier training center focused on transforming student ideas into real-world IT and engineering deliverables. This project establishes a premium online presence and centralized inquiries database."
    p_intro.font.size = Pt(18)
    p_intro.font.color.rgb = WHITE_TEXT
    p_intro.space_after = Pt(24)

    bullets = [
        ("Modern Gen-Z Design System", "Vibrant neon-indigo layout featuring glassmorphism, responsive grids, and fluid scroll-reveal animations."),
        ("Interactive Skills & Portfolios", "Interactive horizontal tabs showcasing training course curricula, vertical ECE/CSE project tabs, and stats counters."),
        ("Persistent Local Database", "A lightweight, custom Node.js server logging queries to a local JSON file without bulky database installations."),
        ("Guarded Administrator Portal", "A restricted dashboard protected by authentication tokens, allowing staff to securely manage client leads.")
    ]

    for title, desc in bullets:
        p = tf2.add_paragraph()
        p.text = f"•  {title}: "
        p.font.bold = True
        p.font.size = Pt(15)
        p.font.color.rgb = CYAN_PRIMARY
        p.space_after = Pt(4)
        
        # Add desc text in same paragraph
        run = p.add_run()
        run.text = desc
        run.font.bold = False
        run.font.color.rgb = MUTED_TEXT
        p.space_after = Pt(12)

    # ----------------------------------------------------
    # SLIDE 3: Visual Design & Frontend Features
    # ----------------------------------------------------
    slide3 = prs.slides.add_slide(slide_layout)
    set_dark_background(slide3)
    add_slide_header(slide3, "Modern Frontend & UI Components")

    # Split into 3 visual cards representing sections
    card_data = [
        ("Hero Terminal", "A simulated shell interface rendering dynamic typing animations (`classify_ai.py`, `wp-config.php`). Boosts professional coder aesthetics instantly."),
        ("Course Tab Matrices", "Smooth horizontal layout switching between Web, Python & AI, Cloud & DevOps, and ECE/CSE branches. Renders custom charts and graphs."),
        ("Stats & Counters", "An Intersection Observer script triggers progress bars to fill up and stats (400+ projects completed) to count up when scrolled into view.")
    ]

    for i, (title, desc) in enumerate(card_data):
        left = Inches(0.8 + i * 4.0)
        card = slide3.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, Inches(2.2), Inches(3.6), Inches(4.2))
        card.fill.solid()
        card.fill.fore_color.rgb = CARD_BG
        card.line.color.rgb = CYAN_PRIMARY
        card.line.width = Pt(1)
        
        tf = card.text_frame
        tf.word_wrap = True
        tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = Inches(0.3)
        
        p = tf.paragraphs[0]
        p.text = title
        p.font.size = Pt(20)
        p.font.bold = True
        p.font.color.rgb = CYAN_PRIMARY
        p.space_after = Pt(16)
        
        p2 = tf.add_paragraph()
        p2.text = desc
        p2.font.size = Pt(14)
        p2.font.color.rgb = MUTED_TEXT
        p2.line_spacing = 1.3

    # ----------------------------------------------------
    # SLIDE 4: Node.js Backend Architecture
    # ----------------------------------------------------
    slide4 = prs.slides.add_slide(slide_layout)
    set_dark_background(slide4)
    add_slide_header(slide4, "Lightweight Node.js Backend")

    # Create two columns
    # Col 1: Details
    tx_col1 = slide4.shapes.add_textbox(Inches(0.8), Inches(2.2), Inches(5.6), Inches(4.5))
    tf_c1 = tx_col1.text_frame
    tf_c1.word_wrap = True
    
    p = tf_c1.paragraphs[0]
    p.text = "Core Backend Features"
    p.font.size = Pt(22)
    p.font.bold = True
    p.font.color.rgb = WHITE_TEXT
    p.space_after = Pt(16)

    features = [
        ("Zero external npm packages", "Runs entirely using built-in modules (`http`, `fs`, `path`) for maximum deployment speed and zero vulnerability."),
        ("Dynamic Port Allocation", "Checks for `EADDRINUSE` errors and automatically scans subsequent ports (starting from 3030). Ensures 100% startup reliability."),
        ("API endpoints routing", "Handles `POST /api/contact` requests, validates forms, and saves data to `contacts.json` locally.")
    ]

    for title, desc in features:
        p = tf_c1.add_paragraph()
        p.text = f"✔  {title}\n"
        p.font.bold = True
        p.font.size = Pt(15)
        p.font.color.rgb = CYAN_PRIMARY
        
        run = p.add_run()
        run.text = desc
        run.font.bold = False
        run.font.size = Pt(13)
        run.font.color.rgb = MUTED_TEXT
        p.space_after = Pt(12)

    # Col 2: Code Editor representation
    code_card = slide4.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(6.8), Inches(2.2), Inches(5.7), Inches(4.2))
    code_card.fill.solid()
    code_card.fill.fore_color.rgb = CARD_BG
    code_card.line.color.rgb = MAGENTA_SECONDARY
    
    tf_code = code_card.text_frame
    tf_code.margin_left = tf_code.margin_right = tf_code.margin_top = tf_code.margin_bottom = Inches(0.4)
    p_code = tf_code.paragraphs[0]
    p_code.text = "server.js snippet:"
    p_code.font.name = 'Courier New'
    p_code.font.size = Pt(14)
    p_code.font.bold = True
    p_code.font.color.rgb = MUTED_TEXT
    p_code.space_after = Pt(12)

    code_lines = [
        "const http = require('http');",
        "const fs = require('fs');",
        "",
        "// Save submission to database",
        "fs.writeFile(DATABASE, JSON.stringify(data), err => {",
        "  if (!err) {",
        "    res.writeHead(200, {'Content-Type': 'json'});",
        "    res.end(JSON.stringify({ success: true }));",
        "  }",
        "});"
    ]
    for line in code_lines:
        p = tf_code.add_paragraph()
        p.text = line
        p.font.name = 'Courier New'
        p.font.size = Pt(13)
        p.font.color.rgb = WHITE_TEXT

    # ----------------------------------------------------
    # SLIDE 5: Security & Admin Dashboard
    # ----------------------------------------------------
    slide5 = prs.slides.add_slide(slide_layout)
    set_dark_background(slide5)
    add_slide_header(slide5, "Security Framework & Admin Panel")

    # Column Left: Security details
    tx_sec = slide5.shapes.add_textbox(Inches(0.8), Inches(2.2), Inches(5.6), Inches(4.5))
    tf_sec = tx_sec.text_frame
    tf_sec.word_wrap = True

    p = tf_sec.paragraphs[0]
    p.text = "Guarding Inquiries & Data"
    p.font.size = Pt(22)
    p.font.bold = True
    p.font.color.rgb = WHITE_TEXT
    p.space_after = Pt(16)

    sec_points = [
        ("Guarded API Endpoint", "The server blocks `GET /api/contacts` requests and throws a `401 Unauthorized` response unless a correct authorization header is supplied."),
        ("Token-based Authorization", "Logs admin credentials on `login.html` via `POST /api/login`. Stores token safely in client's `localStorage` for session authentication."),
        ("Automatic Route Protection", "If the dashboard detects no token, or if the server returns 401, the client automatically wipes local storage and redirects to the login screen.")
    ]

    for title, desc in sec_points:
        p = tf_sec.add_paragraph()
        p.text = f"🔑  {title}\n"
        p.font.bold = True
        p.font.size = Pt(15)
        p.font.color.rgb = MAGENTA_SECONDARY
        
        run = p.add_run()
        run.text = desc
        run.font.bold = False
        run.font.size = Pt(13)
        run.font.color.rgb = MUTED_TEXT
        p.space_after = Pt(14)

    # Column Right: Admin visual layout mockup representation
    admin_card = slide5.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(6.8), Inches(2.2), Inches(5.7), Inches(4.2))
    admin_card.fill.solid()
    admin_card.fill.fore_color.rgb = CARD_BG
    admin_card.line.color.rgb = CYAN_PRIMARY

    tf_admin = admin_card.text_frame
    tf_admin.margin_left = tf_admin.margin_right = tf_admin.margin_top = tf_admin.margin_bottom = Inches(0.4)
    
    p = tf_admin.paragraphs[0]
    p.text = "Dashboard Mockup"
    p.font.size = Pt(18)
    p.font.bold = True
    p.font.color.rgb = CYAN_PRIMARY
    p.space_after = Pt(20)

    rows = [
        ("Prisha Verma", "p.verma@gmail.com", "Web Programming"),
        ("Haley Whitley", "haley@outlook.com", "Academic Kit ECE"),
        ("Amit Kumar", "amit.k@gmail.com", "Python & AI Course")
    ]
    for name, email, query in rows:
        p = tf_admin.add_paragraph()
        p.text = f"👤  {name}   ({email})\n"
        p.font.bold = True
        p.font.size = Pt(13)
        p.font.color.rgb = WHITE_TEXT
        
        run = p.add_run()
        run.text = f"     Inquiry: \"{query}\""
        run.font.bold = False
        run.font.size = Pt(12)
        run.font.color.rgb = MUTED_TEXT
        p.space_after = Pt(10)

    # ----------------------------------------------------
    # SLIDE 6: Deployment Guide (AWS EC2 & Render)
    # ----------------------------------------------------
    slide6 = prs.slides.add_slide(slide_layout)
    set_dark_background(slide6)
    add_slide_header(slide6, "Cloud Deployment Guide")

    # Column Left: AWS EC2 (Traditional VM hosting)
    card_ec2 = slide6.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(2.2), Inches(5.6), Inches(4.5))
    card_ec2.fill.solid()
    card_ec2.fill.fore_color.rgb = CARD_BG
    card_ec2.line.color.rgb = CYAN_PRIMARY

    tf_ec2 = card_ec2.text_frame
    tf_ec2.margin_left = tf_ec2.margin_right = tf_ec2.margin_top = tf_ec2.margin_bottom = Inches(0.4)
    
    p = tf_ec2.paragraphs[0]
    p.text = "Hosting on AWS EC2"
    p.font.size = Pt(20)
    p.font.bold = True
    p.font.color.rgb = CYAN_PRIMARY
    p.space_after = Pt(14)

    steps_ec2 = [
        "1. Launch Ubuntu EC2 Instance.",
        "2. Configure Security Group (Open Ports 80 & 22).",
        "3. SSH into instance and install Node.js & Git.",
        "4. Git clone the repository in the instance.",
        "5. Start application in background on Port 80:\n    'PORT=80 sudo -E pm2 start server.js'"
    ]
    for step in steps_ec2:
        p = tf_ec2.add_paragraph()
        p.text = step
        p.font.size = Pt(13)
        p.font.color.rgb = WHITE_TEXT
        p.space_after = Pt(8)

    # Column Right: Render (Modern serverless container hosting)
    card_render = slide6.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(6.8), Inches(2.2), Inches(5.7), Inches(4.5))
    card_render.fill.solid()
    card_render.fill.fore_color.rgb = CARD_BG
    card_render.line.color.rgb = MAGENTA_SECONDARY

    tf_render = card_render.text_frame
    tf_render.margin_left = tf_render.margin_right = tf_render.margin_top = tf_render.margin_bottom = Inches(0.4)

    p = tf_render.paragraphs[0]
    p.text = "Hosting on Render.com"
    p.font.size = Pt(20)
    p.font.bold = True
    p.font.color.rgb = MAGENTA_SECONDARY
    p.space_after = Pt(14)

    steps_render = [
        "1. Create free account on Render.com.",
        "2. Click 'New Web Service' & connect GitHub.",
        "3. Select repository 'akki07-crypto/sachin'.",
        "4. Configure root directory to 'sachin'.",
        "5. Set start command to 'node server.js'.",
        "6. Click Deploy. Render allocates dynamic port & publishes live URL."
    ]
    for step in steps_render:
        p = tf_render.add_paragraph()
        p.text = step
        p.font.size = Pt(13)
        p.font.color.rgb = WHITE_TEXT
        p.space_after = Pt(8)

    # ----------------------------------------------------
    # SLIDE 7: Future Scope & Conclusion
    # ----------------------------------------------------
    slide7 = prs.slides.add_slide(slide_layout)
    set_dark_background(slide7)
    add_slide_header(slide7, "Future Enhancements & Scope")

    tx_future = slide7.shapes.add_textbox(Inches(0.8), Inches(2.2), Inches(11.7), Inches(4.5))
    tf_f = tx_future.text_frame
    tf_f.word_wrap = True

    f_points = [
        ("Cloud Database Integration", "Transition from a local JSON file to a fully scalable cloud database cluster (e.g. MongoDB Atlas or PostgreSQL) for persistent, secure multi-user concurrent logs."),
        ("Multi-Role User Dashboard", "Implement detailed student login dashboards to download learning tutorials, upload project reports, and review trainer feedback directly on-site."),
        ("Payment Gateway Integration", "Connect secure payment APIs (such as Razorpay or Stripe) to allow direct online registration and fee collection for training modules."),
        ("SMS & Email Notifications", "Automate email/SMS responses confirming inquiry receipt to students, and notify administration teams instantly upon form submission.")
    ]

    for title, desc in f_points:
        p = tf_f.add_paragraph()
        p.text = f"🚀  {title}: "
        p.font.bold = True
        p.font.size = Pt(15)
        p.font.color.rgb = CYAN_PRIMARY
        
        run = p.add_run()
        run.text = desc
        run.font.bold = False
        run.font.color.rgb = MUTED_TEXT
        p.space_after = Pt(16)

    # Save presentation
    prs.save("Tech_Point_Presentation.pptx")
    print("Presentation created successfully as 'Tech_Point_Presentation.pptx'.")

if __name__ == "__main__":
    create_presentation()
